"""GDPval evaluation: rubric-based LLM-as-judge.

GDPval tasks have no reference answers - each task ships a ``rubric_json``
listing criteria with point values. We ask the judge model to score the
candidate answer against each criterion, then aggregate:

- A criterion is satisfied iff the judge returns "yes".
- Score = (achieved points) / (max points).
- ``is_correct`` is True iff every ``required: true`` criterion is satisfied
  AND the achieved score is at least the configured pass threshold (default
  0.5).
- If the rubric is malformed or no criteria are present, returns ``None``
  (unscorable).
"""

from __future__ import annotations

import base64
import json
import logging
import os
import re
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..core.registry import EvaluationRegistry
from .base import EvaluationHandler

LOGGER = logging.getLogger(__name__)

_PASS_THRESHOLD = 0.5
_OUTPUTS_BUDGET_CHARS = 12_000  # per-file budget when summarising deliverables
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
_OFFICE_EXTS = {".xlsx", ".xls", ".xlsm", ".docx", ".pptx"}


def _describe_image_with_vision(path: Path, client) -> str:
    """One-shot image describer using a vision-capable judge model.

    Encodes the image as base64 and asks the judge to produce a detailed
    factual description. The result is reused across every criterion call
    for this deliverable, so cost is O(1 vision call per image), not
    O(n_criteria).
    """
    if not hasattr(client, "_chat_completion"):
        return f"<vision unavailable: client lacks _chat_completion>"
    try:
        with path.open("rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
    except Exception as exc:
        return f"<image read error: {exc}>"
    suf = path.suffix.lower().lstrip(".")
    if suf == "jpg": suf = "jpeg"
    data_url = f"data:image/{suf};base64,{b64}"
    messages = [
        {"role": "system", "content": "You are a careful visual describer."},
        {"role": "user", "content": [
            {"type": "text", "text": (
                "Describe this image factually in 4-8 sentences. Cover: "
                "subject, layout/composition, any visible text or labels, "
                "colour palette, and notable details. Be specific and "
                "literal - no speculation about purpose."
            )},
            {"type": "image_url", "image_url": {"url": data_url}},
        ]},
    ]
    try:
        resp = client._chat_completion(client.model, "", messages=messages)
        return (resp.content or "").strip() or "<empty vision response>"
    except Exception as exc:
        return f"<vision call failed: {exc}>"


def _render_office_to_pdf(src: Path) -> Optional[Path]:
    """Render .docx/.xlsx/.pptx to .pdf using libreoffice headless.

    Returns path to the generated PDF, or None on failure. The PDF
    captures charts, conditional formatting, page layout, embedded
    images etc. that python-docx / openpyxl / python-pptx skip.
    """
    if not src.exists() or src.suffix.lower() not in _OFFICE_EXTS:
        return None
    if not (Path("/usr/bin/soffice").exists() or Path("/usr/bin/libreoffice").exists()):
        return None
    out_dir = Path(tempfile.mkdtemp(prefix="lo_render_"))
    try:
        subprocess.run(
            ["soffice", "--headless", "--norestore", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(src)],
            check=False, capture_output=True, timeout=90,
        )
    except Exception:
        return None
    pdf = out_dir / (src.stem + ".pdf")
    return pdf if pdf.exists() else None

_JUDGE_PROMPT = """You are grading a candidate response against a single rubric criterion.

## Task given to the candidate
{problem}

## Candidate response (plain-text answer)
{response}
{deliverables_section}
## Rubric criterion
{criterion}

Does the candidate response satisfy this criterion?

Reply with exactly one line in this format:
verdict: <yes or no>
reason: <one short sentence>
"""


def _extract_deliverable_text(outputs_dir: Path, vision_client: Any = None) -> str:
    """Return a concatenated text view of any deliverable files.

    Used so the judge can verify rubric criteria that reference produced
    files (e.g. ``Workbook contains tab 'Sample'``).

    Enrichments:

    - **LibreOffice rendering**: for .xlsx/.docx/.pptx, additionally render
      to PDF via headless soffice and extract that text; catches charts,
      conditional formatting, page layout, and embedded images that the
      Python office libraries silently drop.
    - **Vision-based image description**: when ``vision_client`` is
      supplied (a client whose model accepts image_url content), every
      image deliverable gets a 4-8 sentence factual description appended
      to its text section.
    """
    if not outputs_dir or not Path(outputs_dir).is_dir():
        return ""
    sections: list[str] = []
    for path in sorted(Path(outputs_dir).iterdir()):
        if not path.is_file():
            continue
        suf = path.suffix.lower()
        try:
            if suf in (".txt", ".md", ".csv", ".tsv", ".json"):
                text = path.read_text(encoding="utf-8", errors="ignore")
            elif suf in (".xlsx", ".xls", ".xlsm"):
                try:
                    from openpyxl import load_workbook
                except ImportError:
                    text = "<openpyxl not installed>"
                else:
                    wb = load_workbook(filename=str(path), data_only=True, read_only=True)
                    parts: list[str] = []
                    for sheet in wb.worksheets:
                        parts.append(f"# Sheet: {sheet.title}")
                        for row in sheet.iter_rows(values_only=True):
                            cells = [("" if v is None else str(v)) for v in row]
                            parts.append("\t".join(cells))
                    text = "\n".join(parts)
                # Also render to PDF via libreoffice to catch charts /
                # conditional formatting / page layout that openpyxl drops.
                pdf = _render_office_to_pdf(path)
                if pdf is not None:
                    try:
                        import pdfplumber
                        pages = []
                        with pdfplumber.open(str(pdf)) as p:
                            for page in p.pages:
                                pages.append(page.extract_text() or "")
                        rendered = "\n\n".join(pages).strip()
                        if rendered:
                            text += "\n\n## Rendered (libreoffice PDF):\n" + rendered
                    except Exception:
                        pass
            elif suf == ".pdf":
                try:
                    import pdfplumber
                except ImportError:
                    text = "<pdfplumber not installed>"
                else:
                    pages = []
                    with pdfplumber.open(str(path)) as pdf:
                        for page in pdf.pages:
                            pages.append(page.extract_text() or "")
                    text = "\n\n".join(pages)
            elif suf == ".docx":
                try:
                    import docx
                except ImportError:
                    text = "<python-docx not installed>"
                else:
                    text = "\n".join(p.text for p in docx.Document(str(path)).paragraphs)
                pdf = _render_office_to_pdf(path)
                if pdf is not None:
                    try:
                        import pdfplumber
                        with pdfplumber.open(str(pdf)) as p:
                            rendered = "\n\n".join(
                                page.extract_text() or "" for page in p.pages
                            ).strip()
                        if rendered:
                            text += "\n\n## Rendered (libreoffice PDF):\n" + rendered
                    except Exception:
                        pass
            elif suf == ".pptx":
                try:
                    from pptx import Presentation
                except ImportError:
                    text = "<python-pptx not installed>"
                else:
                    prs = Presentation(str(path))
                    parts = []
                    for i, slide in enumerate(prs.slides):
                        parts.append(f"## Slide {i + 1}")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                parts.append(shape.text)
                    text = "\n".join(parts)
                pdf = _render_office_to_pdf(path)
                if pdf is not None:
                    try:
                        import pdfplumber
                        with pdfplumber.open(str(pdf)) as p:
                            rendered = "\n\n".join(
                                page.extract_text() or "" for page in p.pages
                            ).strip()
                        if rendered:
                            text += "\n\n## Rendered (libreoffice PDF):\n" + rendered
                    except Exception:
                        pass
            elif suf == ".ipynb":
                try:
                    nb = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
                    parts = []
                    for i, cell in enumerate(nb.get("cells") or []):
                        src = cell.get("source")
                        if isinstance(src, list):
                            src = "".join(src)
                        parts.append(f"## Cell {i + 1} ({cell.get('cell_type', '?')})\n{src or ''}")
                    text = "\n\n".join(parts)
                except Exception as exc:
                    text = f"<ipynb parse error: {exc}>"
            elif suf == ".zip":
                try:
                    import zipfile
                    with zipfile.ZipFile(str(path)) as z:
                        names = z.namelist()
                    text = "Archive contents:\n" + "\n".join(
                        f"  {n}" for n in names[:200]
                    ) + (f"\n  ... [{len(names) - 200} more]" if len(names) > 200 else "")
                except Exception as exc:
                    text = f"<zip read error: {exc}>"
            elif suf in (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"):
                try:
                    from PIL import Image
                    img = Image.open(str(path))
                    text = f"Image file. format={img.format} mode={img.mode} size={img.size}"
                except Exception as exc:
                    text = f"<image read error: {exc}>"
                # When a vision-capable client is available, ask it to
                # describe the image so the rubric judge can evaluate visual
                # criteria (composition, labels, color, etc.) using text.
                if vision_client is not None:
                    desc = _describe_image_with_vision(path, vision_client)
                    if desc:
                        text += "\n\n## Vision description:\n" + desc
            elif suf in (".wav", ".mp3", ".m4a", ".flac", ".ogg"):
                # Transcribe with whisper (~base model). Limited to first
                # ~30 s of audio to keep grade time bounded.
                try:
                    import whisper
                    model = whisper.load_model("base")
                    result = model.transcribe(str(path), verbose=False)
                    transcript = (result.get("text") or "")[:_OUTPUTS_BUDGET_CHARS]
                    text = f"Audio transcript:\n{transcript}"
                except Exception as exc:
                    text = f"<audio transcribe error: {exc}>"
            elif suf == ".mp4":
                # Extract audio with ffmpeg, transcribe with whisper
                import subprocess
                import tempfile
                wav = path.with_suffix(".wav")
                try:
                    if not wav.exists():
                        wav = Path(tempfile.NamedTemporaryFile(
                            delete=False, suffix=".wav"
                        ).name)
                        subprocess.run(
                            ["ffmpeg", "-y", "-i", str(path), "-vn", "-ar", "16000",
                             "-ac", "1", str(wav)],
                            capture_output=True, check=True, timeout=120,
                        )
                    import whisper
                    model = whisper.load_model("base")
                    result = model.transcribe(str(wav), verbose=False)
                    transcript = (result.get("text") or "")[:_OUTPUTS_BUDGET_CHARS]
                    text = f"Video file (audio transcribed):\n{transcript}"
                except Exception as exc:
                    text = f"<video transcribe error: {exc}>"
            elif suf == ".psd":
                try:
                    from psd_tools import PSDImage
                    psd = PSDImage.open(str(path))
                    text = (
                        f"PSD file. size={psd.size} layers="
                        + ",".join(layer.name for layer in psd.descendants() if layer.name)[:1000]
                    )
                except Exception as exc:
                    text = f"<psd read error: {exc}>"
            elif suf == ".step":
                # CAD file. No installable reader on this box; describe only.
                text = f"STEP CAD file. size={path.stat().st_size}B  (binary CAD content not parsed)"
            elif suf == ".py":
                # Treat py deliverables like text (judge will read the source)
                text = path.read_text(encoding="utf-8", errors="ignore")
            else:
                text = path.read_text(encoding="utf-8", errors="ignore")
        except Exception as exc:
            text = f"<error reading {path.name}: {exc}>"

        if len(text) > _OUTPUTS_BUDGET_CHARS:
            text = text[:_OUTPUTS_BUDGET_CHARS] + "\n... [truncated]"
        sections.append(f"### Deliverable file: {path.name}\n{text}")
    return "\n\n".join(sections)


def _parse_rubric(raw: Any) -> List[Dict[str, Any]]:
    """Return the list of rubric items from a rubric_json blob."""
    if not raw:
        return []
    if isinstance(raw, list):
        items = raw
    elif isinstance(raw, str):
        try:
            parsed = json.loads(raw)
        except json.JSONDecodeError:
            return []
        if isinstance(parsed, list):
            items = parsed
        elif isinstance(parsed, dict):
            # Some rubrics wrap the list under a key
            items = (
                parsed.get("rubric")
                or parsed.get("criteria")
                or parsed.get("items")
                or []
            )
        else:
            return []
    else:
        return []

    cleaned: list[dict[str, Any]] = []
    for item in items:
        if isinstance(item, dict) and item.get("criterion"):
            cleaned.append(item)
    return cleaned


def _judge_verdict(raw: str) -> bool:
    """Return True iff the judge said 'verdict: yes'."""
    if not raw:
        return False
    m = re.search(r"verdict\s*:\s*(yes|no)", raw, re.IGNORECASE)
    if m:
        return m.group(1).lower() == "yes"
    # Fallback: any standalone yes/no token in the first line
    first_line = raw.strip().splitlines()[0] if raw.strip() else ""
    return bool(re.search(r"\byes\b", first_line, re.IGNORECASE))


@EvaluationRegistry.register("gdpval")
class GdpvalHandler(EvaluationHandler):
    """Rubric-based LLM-as-judge grading for GDPval tasks."""

    evaluation_method = "gdpval"

    def evaluate(
        self,
        *,
        problem: str,
        reference: str,
        model_answer: str,
        metadata: Dict[str, object],
    ) -> Tuple[Optional[bool], Dict[str, object]]:
        if not model_answer or not model_answer.strip():
            return False, {"reason": "empty_response"}

        rubric_raw = metadata.get("rubric_json")
        rubric = _parse_rubric(rubric_raw)
        if not rubric:
            return None, {"reason": "no_rubric"}

        if not hasattr(self._client, "chat"):
            return None, {"reason": "no_llm_client_for_judging"}

        # Optional deliverable-file content: the file-io agent writes outputs
        # under <workspace>/outputs/ and surfaces the path via metadata.
        outputs_dir = metadata.get("gdpval_outputs_dir")
        if not outputs_dir:
            # Conventional path: workspace sibling of metadata's instance_id
            # - fall back to scanning common locations.
            pass
        deliverable_text = (
            _extract_deliverable_text(Path(outputs_dir), vision_client=self._client)
            if outputs_dir else ""
        )
        deliverables_section = (
            f"\n## Deliverable files produced\n{deliverable_text}\n"
            if deliverable_text
            else ""
        )

        per_criterion: list[dict[str, Any]] = []
        achieved_points = 0.0
        max_points = 0.0
        required_failures = 0

        for idx, item in enumerate(rubric):
            criterion_text = str(item.get("criterion") or "").strip()
            if not criterion_text:
                continue

            try:
                points = float(item.get("score", 1.0))
            except (TypeError, ValueError):
                points = 1.0
            required = bool(item.get("required", False))
            rubric_item_id = item.get("rubric_item_id", f"item_{idx}")

            max_points += points

            prompt = _JUDGE_PROMPT.format(
                problem=problem[:4000],  # keep judge prompt bounded
                response=model_answer[:8000],
                deliverables_section=deliverables_section,
                criterion=criterion_text,
            )
            try:
                raw = self._client.chat(
                    system_prompt="",
                    user_prompt=prompt,
                    temperature=0.0,
                    max_output_tokens=512,
                )
            except Exception as exc:
                LOGGER.warning(
                    "Rubric judge failed for item %s: %s", rubric_item_id, exc
                )
                per_criterion.append(
                    {
                        "rubric_item_id": rubric_item_id,
                        "criterion": criterion_text,
                        "points": points,
                        "required": required,
                        "satisfied": False,
                        "error": str(exc),
                    }
                )
                if required:
                    required_failures += 1
                continue

            satisfied = _judge_verdict(raw)
            if satisfied:
                achieved_points += points
            elif required:
                required_failures += 1

            per_criterion.append(
                {
                    "rubric_item_id": rubric_item_id,
                    "criterion": criterion_text,
                    "points": points,
                    "required": required,
                    "satisfied": satisfied,
                    "judge_output": raw[:500],
                }
            )

        score = achieved_points / max_points if max_points > 0 else 0.0
        is_correct = required_failures == 0 and score >= _PASS_THRESHOLD

        return is_correct, {
            "match_type": "rubric_llm_judge",
            "score": score,
            "achieved_points": achieved_points,
            "max_points": max_points,
            "required_failures": required_failures,
            "num_criteria": len(per_criterion),
            "per_criterion": per_criterion,
            "pass_threshold": _PASS_THRESHOLD,
        }


__all__ = ["GdpvalHandler"]
