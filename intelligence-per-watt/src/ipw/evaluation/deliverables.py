"""Deliverable-aware evaluation helpers.

This module converts files produced by agentic tasks into bounded,
judge-readable evidence. It intentionally avoids assuming a single output
format: GDPval tasks can ask for documents, spreadsheets, images, audio,
video, archives, or mixed bundles.
"""

from __future__ import annotations

import hashlib
import json
import logging
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable

LOGGER = logging.getLogger(__name__)

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".xml",
    ".html",
    ".htm",
    ".svg",
    ".py",
    ".js",
    ".ts",
    ".css",
    ".yaml",
    ".yml",
    ".log",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".flac", ".ogg", ".aac"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".mkv", ".avi", ".webm"}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xlsm", ".xltx", ".xltm"}


@dataclass
class FileEvidence:
    """A bounded, text-first description of a deliverable or reference file."""

    path: str
    name: str
    suffix: str
    mime_type: str | None
    size_bytes: int
    sha256: str
    kind: str
    text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    children: list["FileEvidence"] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        data = asdict(self)
        data["children"] = [child.to_dict() for child in self.children]
        return data

    def to_prompt_section(self, *, max_chars: int = 16_000, depth: int = 0) -> str:
        indent = "  " * depth
        meta = json.dumps(self.metadata, default=str, ensure_ascii=False)[:4000]
        text = self.text
        if len(text) > max_chars:
            text = text[:max_chars] + "\n... [truncated]"
        lines = [
            f"{indent}### File: {self.name}",
            f"{indent}- path: {self.path}",
            f"{indent}- kind: {self.kind}",
            f"{indent}- mime_type: {self.mime_type}",
            f"{indent}- size_bytes: {self.size_bytes}",
            f"{indent}- sha256: {self.sha256}",
            f"{indent}- metadata: {meta}",
        ]
        if self.errors:
            lines.append(f"{indent}- extraction_errors: {json.dumps(self.errors, ensure_ascii=False)}")
        if text:
            lines.append(f"{indent}Content/evidence:\n{text}")
        for child in self.children:
            lines.append(child.to_prompt_section(max_chars=max_chars // 2, depth=depth + 1))
        return "\n".join(lines)


def _sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _bounded(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n... [truncated]"


def _safe_read_text(path: Path, max_chars: int) -> str:
    return _bounded(path.read_text(encoding="utf-8", errors="ignore"), max_chars)


def _inspect_text(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    return _safe_read_text(path, max_chars), {}, []


def _inspect_spreadsheet(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    errors: list[str] = []
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "", {}, ["openpyxl is not installed"]

    parts: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        wb_values = load_workbook(str(path), data_only=True, read_only=True)
        wb_formulas = load_workbook(str(path), data_only=False, read_only=True)
        metadata["sheet_names"] = wb_values.sheetnames
        for ws_values, ws_formulas in zip(wb_values.worksheets, wb_formulas.worksheets):
            parts.append(f"# Sheet: {ws_values.title}")
            metadata.setdefault("sheets", []).append(
                {
                    "name": ws_values.title,
                    "max_row": ws_values.max_row,
                    "max_column": ws_values.max_column,
                }
            )
            row_limit = min(ws_values.max_row or 0, 80)
            col_limit = min(ws_values.max_column or 0, 30)
            for r_idx, row in enumerate(
                ws_values.iter_rows(min_row=1, max_row=row_limit, max_col=col_limit, values_only=True),
                start=1,
            ):
                formula_row = next(
                    ws_formulas.iter_rows(min_row=r_idx, max_row=r_idx, max_col=col_limit, values_only=True)
                )
                cells = []
                for value, formula in zip(row, formula_row):
                    if formula is not None and isinstance(formula, str) and formula.startswith("="):
                        cells.append(formula)
                    else:
                        cells.append("" if value is None else str(value))
                if any(cells):
                    parts.append("\t".join(cells))
            if (ws_values.max_row or 0) > row_limit:
                parts.append(f"... [{(ws_values.max_row or 0) - row_limit} more rows]")
        return _bounded("\n".join(parts), max_chars), metadata, errors
    except Exception as exc:
        return "", metadata, [f"spreadsheet parse error: {exc}"]


def _inspect_docx(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    try:
        import docx
    except ImportError:
        return "", {}, ["python-docx is not installed"]
    try:
        doc = docx.Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for t_idx, table in enumerate(doc.tables, start=1):
            parts.append(f"# Table {t_idx}")
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return _bounded("\n".join(parts), max_chars), {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}, []
    except Exception as exc:
        return "", {}, [f"docx parse error: {exc}"]


def _inspect_pptx(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    try:
        from pptx import Presentation
    except ImportError:
        return "", {}, ["python-pptx is not installed"]
    try:
        prs = Presentation(str(path))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {idx}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return _bounded("\n".join(parts), max_chars), {"slides": len(prs.slides)}, []
    except Exception as exc:
        return "", {}, [f"pptx parse error: {exc}"]


def _inspect_pdf(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    errors: list[str] = []
    parts: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            metadata["pages"] = len(pdf.pages)
            for idx, page in enumerate(pdf.pages[:20], start=1):
                parts.append(f"# Page {idx}")
                parts.append(page.extract_text() or "")
            if len(pdf.pages) > 20:
                parts.append(f"... [{len(pdf.pages) - 20} more pages]")
    except ImportError:
        errors.append("pdfplumber is not installed")
    except Exception as exc:
        errors.append(f"pdf text parse error: {exc}")

    if not "".join(parts).strip():
        try:
            result = subprocess.run(
                ["pdftotext", str(path), "-"],
                capture_output=True,
                text=True,
                timeout=60,
                check=False,
            )
            if result.stdout:
                parts.append(result.stdout)
            if result.stderr:
                errors.append(result.stderr[:500])
        except Exception as exc:
            errors.append(f"pdftotext failed: {exc}")
    return _bounded("\n".join(parts), max_chars), metadata, errors


def _inspect_image(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    errors: list[str] = []
    metadata: dict[str, Any] = {}
    text_parts: list[str] = []
    try:
        from PIL import Image

        with Image.open(str(path)) as img:
            metadata.update(
                {
                    "format": img.format,
                    "mode": img.mode,
                    "width": img.width,
                    "height": img.height,
                }
            )
    except Exception as exc:
        errors.append(f"image metadata error: {exc}")

    if shutil.which("tesseract"):
        try:
            result = subprocess.run(
                ["tesseract", str(path), "stdout"],
                capture_output=True,
                text=True,
                timeout=90,
                check=False,
            )
            if result.stdout.strip():
                text_parts.append("# OCR text")
                text_parts.append(result.stdout)
            if result.stderr and result.returncode != 0:
                errors.append(result.stderr[:500])
        except Exception as exc:
            errors.append(f"image OCR failed: {exc}")
    else:
        errors.append("tesseract not available; visual content represented by metadata only")

    return _bounded("\n".join(text_parts), max_chars), metadata, errors


def _inspect_audio(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    errors: list[str] = []
    metadata: dict[str, Any] = {}
    text_parts: list[str] = []
    try:
        from mutagen import File as MutagenFile

        audio = MutagenFile(str(path))
        if audio is not None and audio.info is not None:
            metadata["duration_seconds"] = getattr(audio.info, "length", None)
            metadata["bitrate"] = getattr(audio.info, "bitrate", None)
            metadata["sample_rate"] = getattr(audio.info, "sample_rate", None)
            metadata["channels"] = getattr(audio.info, "channels", None)
    except ImportError:
        errors.append("mutagen is not installed")
    except Exception as exc:
        errors.append(f"audio metadata error: {exc}")

    if os.getenv("IPW_ENABLE_MEDIA_TRANSCRIPTION", "").lower() in {"1", "true", "yes"}:
        try:
            import whisper

            model_name = os.getenv("IPW_WHISPER_MODEL", "base")
            model = whisper.load_model(model_name)
            result = model.transcribe(str(path), verbose=False)
            transcript = result.get("text") or ""
            if transcript:
                text_parts.append("# Transcript")
                text_parts.append(transcript)
        except ImportError:
            errors.append("whisper is not installed")
        except Exception as exc:
            errors.append(f"audio transcription failed: {exc}")
    else:
        errors.append("media transcription disabled; set IPW_ENABLE_MEDIA_TRANSCRIPTION=1 to enable")

    return _bounded("\n".join(text_parts), max_chars), metadata, errors


def _ffprobe_metadata(path: Path) -> dict[str, Any]:
    if not shutil.which("ffprobe"):
        return {"ffprobe": "not available"}
    result = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_format",
            "-show_streams",
            "-of",
            "json",
            str(path),
        ],
        capture_output=True,
        text=True,
        timeout=60,
        check=False,
    )
    if not result.stdout:
        return {"ffprobe_error": result.stderr[:500]}
    try:
        return json.loads(result.stdout)
    except json.JSONDecodeError:
        return {"ffprobe_raw": result.stdout[:2000]}


def _inspect_video(path: Path, max_chars: int) -> tuple[str, dict[str, Any], list[str]]:
    errors: list[str] = []
    metadata = _ffprobe_metadata(path)
    text_parts: list[str] = []

    if os.getenv("IPW_ENABLE_MEDIA_TRANSCRIPTION", "").lower() in {"1", "true", "yes"}:
        with tempfile.TemporaryDirectory() as tmp:
            wav = Path(tmp) / "audio.wav"
            try:
                if shutil.which("ffmpeg"):
                    subprocess.run(
                        ["ffmpeg", "-y", "-i", str(path), "-vn", "-ar", "16000", "-ac", "1", str(wav)],
                        capture_output=True,
                        text=True,
                        timeout=180,
                        check=False,
                    )
                    if wav.exists():
                        audio_text, audio_meta, audio_errors = _inspect_audio(wav, max_chars)
                        text_parts.append(audio_text)
                        metadata["extracted_audio"] = audio_meta
                        errors.extend(audio_errors)
                else:
                    errors.append("ffmpeg not available for video audio extraction")
            except Exception as exc:
                errors.append(f"video audio extraction failed: {exc}")
    else:
        errors.append("media transcription disabled; set IPW_ENABLE_MEDIA_TRANSCRIPTION=1 to enable")

    return _bounded("\n".join(t for t in text_parts if t), max_chars), metadata, errors


def inspect_file(path: Path, *, max_chars: int = 16_000, archive_depth: int = 1) -> FileEvidence:
    """Return bounded evidence for a single file."""
    path = Path(path)
    suffix = path.suffix.lower()
    mime_type = mimetypes.guess_type(str(path))[0]
    stat = path.stat()
    evidence = FileEvidence(
        path=str(path),
        name=path.name,
        suffix=suffix,
        mime_type=mime_type,
        size_bytes=stat.st_size,
        sha256=_sha256(path),
        kind="unknown",
    )

    try:
        if suffix in TEXT_EXTENSIONS:
            evidence.kind = "text"
            evidence.text, evidence.metadata, evidence.errors = _inspect_text(path, max_chars)
        elif suffix in SPREADSHEET_EXTENSIONS:
            evidence.kind = "spreadsheet"
            evidence.text, evidence.metadata, evidence.errors = _inspect_spreadsheet(path, max_chars)
        elif suffix == ".docx":
            evidence.kind = "document"
            evidence.text, evidence.metadata, evidence.errors = _inspect_docx(path, max_chars)
        elif suffix == ".pptx":
            evidence.kind = "presentation"
            evidence.text, evidence.metadata, evidence.errors = _inspect_pptx(path, max_chars)
        elif suffix == ".pdf":
            evidence.kind = "pdf"
            evidence.text, evidence.metadata, evidence.errors = _inspect_pdf(path, max_chars)
        elif suffix in IMAGE_EXTENSIONS:
            evidence.kind = "image"
            evidence.text, evidence.metadata, evidence.errors = _inspect_image(path, max_chars)
        elif suffix in AUDIO_EXTENSIONS:
            evidence.kind = "audio"
            evidence.text, evidence.metadata, evidence.errors = _inspect_audio(path, max_chars)
        elif suffix in VIDEO_EXTENSIONS:
            evidence.kind = "video"
            evidence.text, evidence.metadata, evidence.errors = _inspect_video(path, max_chars)
        elif suffix == ".zip":
            evidence.kind = "archive"
            evidence.metadata = {"archive_type": "zip"}
            if archive_depth > 0:
                with tempfile.TemporaryDirectory() as tmp:
                    root = Path(tmp)
                    with zipfile.ZipFile(str(path)) as zf:
                        names = zf.namelist()
                        evidence.metadata["entries"] = names[:200]
                        evidence.metadata["entry_count"] = len(names)
                        zf.extractall(root)
                    for child in sorted(root.rglob("*")):
                        if child.is_file() and len(evidence.children) < 30:
                            evidence.children.append(
                                inspect_file(child, max_chars=max_chars // 2, archive_depth=archive_depth - 1)
                            )
            else:
                with zipfile.ZipFile(str(path)) as zf:
                    names = zf.namelist()
                evidence.text = "\n".join(names[:500])
                evidence.metadata["entry_count"] = len(names)
        else:
            evidence.kind = "binary"
            evidence.errors.append("No specialized parser for this file type; using metadata only")
    except Exception as exc:
        evidence.errors.append(f"file inspection failed: {exc}")

    return evidence


def collect_file_evidence(
    paths: Iterable[str | Path],
    *,
    max_files: int = 60,
    max_chars_per_file: int = 16_000,
) -> list[FileEvidence]:
    """Inspect all existing files in paths, expanding directories recursively."""
    files: list[Path] = []
    for raw in paths:
        if raw is None:
            continue
        path = Path(raw)
        if path.is_dir():
            files.extend(sorted(p for p in path.rglob("*") if p.is_file()))
        elif path.is_file():
            files.append(path)

    evidence: list[FileEvidence] = []
    seen: set[Path] = set()
    for path in files:
        try:
            resolved = path.resolve()
        except Exception:
            resolved = path
        if resolved in seen:
            continue
        seen.add(resolved)
        if len(evidence) >= max_files:
            break
        evidence.append(inspect_file(path, max_chars=max_chars_per_file))
    return evidence


def render_evidence_sections(title: str, evidence: list[FileEvidence], *, max_chars: int = 80_000) -> str:
    """Render evidence as a bounded prompt section."""
    parts = [f"## {title}"]
    for item in evidence:
        parts.append(item.to_prompt_section())
    rendered = "\n\n".join(parts)
    return _bounded(rendered, max_chars)


def extract_json_object(text: str) -> dict[str, Any] | None:
    """Best-effort extraction of a JSON object from a judge response."""
    if not text:
        return None
    stripped = text.strip()

    fenced_blocks = re.findall(r"```(?:json)?\s*(\{.*?\})\s*```", stripped, flags=re.DOTALL)
    for block in fenced_blocks:
        try:
            parsed = json.loads(block)
            if isinstance(parsed, dict):
                return parsed
        except json.JSONDecodeError:
            continue

    if stripped.startswith("```"):
        stripped = re.sub(r"^```(?:json)?\s*", "", stripped)
        stripped = re.sub(r"\s*```$", "", stripped)
    try:
        parsed = json.loads(stripped)
        return parsed if isinstance(parsed, dict) else None
    except json.JSONDecodeError:
        pass

    decoder = json.JSONDecoder()
    for match in re.finditer(r"\{", stripped):
        try:
            parsed, _ = decoder.raw_decode(stripped[match.start():])
        except json.JSONDecodeError:
            continue
        if isinstance(parsed, dict):
            return parsed
    return None


__all__ = [
    "FileEvidence",
    "collect_file_evidence",
    "extract_json_object",
    "inspect_file",
    "render_evidence_sections",
]
